import pptx
import io
import base64
from typing import Dict, Any, List
from .text_normalizer import normalize_text

def process_pptx(file_path: str, extract_images: bool = True) -> Dict[str, Any]:
    """
    Processes PPTX files using python-pptx:
    1. Extracts slide titles, text frames, speaker notes, and tables.
    2. Extracts embedded slide images/charts for visual analysis.
    Returns:
    {
      "text": str,
      "slide_images": [{"slide": int, "type": "slide_image", "base64_image": str}],
      "metadata": {"total_slides": int, "images_count": int}
    }
    """
    prs = pptx.Presentation(file_path)
    slide_blocks = []
    slide_images = []
    
    total_slides = len(prs.slides)
    
    for idx, slide in enumerate(prs.slides):
        slide_num = idx + 1
        lines = [f"--- Slide {slide_num} ---"]
        
        # 1. Slide Title & Shapes Text
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    t = paragraph.text.strip()
                    if t:
                        lines.append(t)
            elif shape.has_table:
                table = shape.table
                table_lines = []
                for i, row in enumerate(table.rows):
                    cells = [cell.text.strip().replace("\n", " ") for cell in row.cells]
                    table_lines.append(" | ".join(cells))
                    if i == 0:
                        table_lines.append(" | ".join(["---"] * len(cells)))
                if table_lines:
                    lines.append("\n" + "\n".join(table_lines))
            elif extract_images and shape.shape_type == pptx.enum.shapes.MSO_SHAPE_TYPE.PICTURE:
                try:
                    image = shape.image
                    img_bytes = image.blob
                    b64_img = base64.b64encode(img_bytes).decode("utf-8")
                    slide_images.append({
                        "slide": slide_num,
                        "type": "slide_image",
                        "base64_image": b64_img
                    })
                except Exception as img_err:
                    print(f"[PPTX Processor] Image extraction warning on slide {slide_num}: {img_err}")
                    
        # 2. Speaker Notes
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                lines.append(f"\n*Speaker Notes:* {notes}")
                
        slide_blocks.append("\n".join(lines))
        
    full_text = normalize_text("\n\n".join(slide_blocks))
    
    return {
        "text": full_text,
        "slide_images": slide_images,
        "metadata": {
            "total_slides": total_slides,
            "images_count": len(slide_images)
        }
    }
